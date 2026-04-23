"""Create Step 5 PPT slides matching CUHK purple theme."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
import os

SCRIPT_DIR = Path(__file__).resolve().parent
REPORT_FIG_DIR = SCRIPT_DIR / "report_figures"
OUTPUT_PPT = Path(r"C:\Users\93480\Desktop\step5_slides.pptx")

# CUHK Purple theme colors
PURPLE_DARK = RGBColor(0x5B, 0x2C, 0x6F)   # dark purple header
PURPLE_MID  = RGBColor(0x7B, 0x2D, 0x8E)   # medium purple
PURPLE_LIGHT = RGBColor(0x9B, 0x59, 0xB6)  # light purple accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GOLD = RGBColor(0xC0, 0x9A, 0x36)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_purple_header(slide, title_text):
    """Add a purple header bar with title, matching the PPT style."""
    # Purple header rectangle
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE_DARK
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.1), Inches(10), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"


def add_bullet_text(text_frame, text, level=0, size=18, bold=False, color=DARK_GRAY):
    """Add a bullet point to a text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_after = Pt(6)
    return p


def add_first_bullet(text_frame, text, size=18, bold=False, color=DARK_GRAY):
    """Set text on the first paragraph."""
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_after = Pt(6)
    return p


# ============================================================
# Slide 1: Section Title - Textual Generate
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
# Full purple background
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = PURPLE_DARK
bg.line.fill.background()

# Section title
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Step 5: (Optional) Textual Generate"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.LEFT

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "Adding Creative Text Captions to Generated Stickers"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0xE0, 0xD0, 0xF0)
p2.font.name = "Calibri"

# ============================================================
# Slide 2: Overview / Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_purple_header(slide, "Overview: Text Overlay Pipeline")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True

add_first_bullet(tf, "Goal: Add short, creative text captions to enhance sticker personality", size=20, bold=True)
add_bullet_text(tf, "Challenge: SD 1.5 cannot generate legible text in images", size=18)
add_bullet_text(tf, "Solution: Post-processing pipeline with VLM + PIL overlay", size=18)
add_bullet_text(tf, "", size=12)
add_bullet_text(tf, "Three-Stage Pipeline:", size=20, bold=True, color=PURPLE_MID)
add_bullet_text(tf, "1. Attempt direct text generation with SD → Failed (garbled text)", size=18, level=1)
add_bullet_text(tf, "2. Use VLM (GPT-4o-mini) to auto-generate captions from sticker images", size=18, level=1)
add_bullet_text(tf, "3. Apply style-aware text rendering with contrast detection", size=18, level=1)

# Pipeline flow boxes
box_y = Inches(5.2)
box_h = Inches(0.8)
box_w = Inches(2.5)
labels = ["Sticker Image", "VLM Caption", "Style-Aware\nOverlay", "Final Sticker"]
colors = [PURPLE_LIGHT, PURPLE_MID, PURPLE_DARK, GOLD]
for i, (label, color) in enumerate(zip(labels, colors)):
    x = Inches(0.8 + i * 3.1)
    box = slide.shapes.add_shape(1, x, box_y, box_w, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.text_frame.paragraphs[0].text = label
    box.text_frame.paragraphs[0].font.size = Pt(14)
    box.text_frame.paragraphs[0].font.color.rgb = WHITE
    box.text_frame.paragraphs[0].font.bold = True
    box.text_frame.paragraphs[0].font.name = "Calibri"
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow between boxes
    if i < 3:
        arrow_x = Inches(0.8 + (i + 1) * 3.1 - 0.5)
        arr = slide.shapes.add_shape(1, arrow_x, Inches(5.4), Inches(0.4), Inches(0.4))
        arr.fill.solid()
        arr.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
        arr.line.fill.background()
        arr.text_frame.paragraphs[0].text = "→"
        arr.text_frame.paragraphs[0].font.size = Pt(20)
        arr.text_frame.paragraphs[0].font.color.rgb = WHITE
        arr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# Slide 3: Initial Attempt - SD Text Generation Failed
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_purple_header(slide, "Initial Attempt: Direct Text Generation with SD")

# Left side: text
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True

add_first_bullet(tf, "What we tried:", size=20, bold=True, color=PURPLE_MID)
add_bullet_text(tf, 'Added text requests in SD prompt:\n"a cat sticker with text \'Cool Cat\'"', size=16, level=1)
add_bullet_text(tf, "", size=8)
add_bullet_text(tf, "Result: Completely illegible text", size=20, bold=True, color=RGBColor(0xCC, 0x00, 0x00))
add_bullet_text(tf, "Garbled characters, distorted letters", size=16, level=1)
add_bullet_text(tf, '"CLAL COOL COL COol" instead of "Cool Cat"', size=16, level=1)
add_bullet_text(tf, "", size=8)
add_bullet_text(tf, "Root Cause:", size=20, bold=True, color=PURPLE_MID)
add_bullet_text(tf, "SD 1.5 lacks dedicated text rendering modules", size=16, level=1)
add_bullet_text(tf, "Known limitation of diffusion models", size=16, level=1)
add_bullet_text(tf, "", size=8)
add_bullet_text(tf, "Decision: Switch to post-processing approach", size=18, bold=True, color=PURPLE_DARK)

# Right side: failed images
# Check if failed examples exist; add placeholder text if not
failed_img_1 = SCRIPT_DIR / "report_figures" / "sd_failed_text_1.png"
failed_img_2 = SCRIPT_DIR / "report_figures" / "sd_failed_text_2.png"

# Add "Failed Examples" label
label_box = slide.shapes.add_textbox(Inches(7.2), Inches(1.2), Inches(5), Inches(0.5))
lp = label_box.text_frame.paragraphs[0]
lp.text = "SD-Generated Text Examples (Failed)"
lp.font.size = Pt(16)
lp.font.bold = True
lp.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
lp.font.name = "Calibri"

if failed_img_1.exists():
    slide.shapes.add_picture(str(failed_img_1), Inches(7.2), Inches(1.8), Inches(2.4), Inches(2.4))
if failed_img_2.exists():
    slide.shapes.add_picture(str(failed_img_2), Inches(10), Inches(1.8), Inches(2.4), Inches(2.4))

if not failed_img_1.exists():
    # Placeholder box
    box = slide.shapes.add_shape(1, Inches(7.2), Inches(1.8), Inches(5), Inches(4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    box.text_frame.paragraphs[0].text = "Insert SD failed text images here\n(sd_failed_text_1.png, sd_failed_text_2.png)"
    box.text_frame.paragraphs[0].font.size = Pt(14)
    box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# ============================================================
# Slide 4: VLM Caption Generation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_purple_header(slide, "VLM-Based Caption Generation")

# Left column
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(6), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True

add_first_bullet(tf, "Model: GPT-4o-mini via POE API", size=18, bold=True)
add_bullet_text(tf, "Automatically analyzes sticker image content", size=16, level=1)
add_bullet_text(tf, "", size=8)
add_bullet_text(tf, "Output Format: Adjective + Noun", size=18, bold=True, color=PURPLE_MID)
add_bullet_text(tf, '"Sleepy Whiskers"  "Brave Bear"  "Chill Penguin"', size=16, level=1)
add_bullet_text(tf, '"Joyful Kangaroo"  "Sassy Furball"  "Gentle Giraffe"', size=16, level=1)
add_bullet_text(tf, "", size=8)
add_bullet_text(tf, "Diversity Constraints:", size=18, bold=True, color=PURPLE_MID)
add_bullet_text(tf, "For cats: rotate among Cat, Kitty, Whiskers,\nFurball, Meowster, Paws, Fluffball", size=16, level=1)
add_bullet_text(tf, "Prevents repetitive captions across sticker set", size=16, level=1)
add_bullet_text(tf, "", size=8)
add_bullet_text(tf, "Caching: Results saved to JSON for reproducibility", size=16, bold=False)

# Right column - example box
box = slide.shapes.add_shape(1, Inches(7.5), Inches(1.3), Inches(5), Inches(5.2))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xFA)
box.line.color.rgb = PURPLE_LIGHT

title_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.5), Inches(4.5), Inches(0.5))
tp = title_box.text_frame.paragraphs[0]
tp.text = "Generated Caption Examples"
tp.font.size = Pt(16)
tp.font.bold = True
tp.font.color.rgb = PURPLE_DARK
tp.font.name = "Calibri"

examples = [
    ("animal_1.png", "Cheerful Paws"),
    ("animal_5.png", "Happy Penguin"),
    ("animal_10.png", "Fluffy Meowster"),
    ("animal_15.png", "Gentle Giraffe"),
    ("animal_20.png", "Curious Meowster"),
    ("animal_24.png", "Cheerful Capybara"),
    ("animal_27.png", "Fluffy Queen"),
    ("animal_30.png", "Chill Furball"),
    ("animal_33.png", "Gentle Tower"),
    ("animal_37.png", "Charming Bear"),
]

ex_box = slide.shapes.add_textbox(Inches(7.8), Inches(2.1), Inches(4.5), Inches(4))
ex_tf = ex_box.text_frame
ex_tf.word_wrap = True
for i, (img, cap) in enumerate(examples):
    if i == 0:
        p = ex_tf.paragraphs[0]
    else:
        p = ex_tf.add_paragraph()
    p.text = f"  {img:20s} →  {cap}"
    p.font.size = Pt(13)
    p.font.name = "Consolas"
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(3)

# ============================================================
# Slide 5: Style-Aware Text Rendering
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_purple_header(slide, "Style-Aware Text Rendering")

# Three technique boxes
techniques = [
    ("Dominant Color\nExtraction", 
     "• Downsample to 50×50px\n• Filter extreme brightness\n• Find most frequent\n  saturated color\n• Use as text color base"),
    ("Local Background\nBrightness Sampling",
     "• Sample region behind\n  text bounding box\n• More accurate than\n  global statistics\n• Detect local contrast"),
    ("Contrast\nAdjustment",
     "• Compare text vs background\n• If difference < 60/255:\n  auto-swap colors\n• Ensures readability on\n  any background"),
]

for i, (title, desc) in enumerate(techniques):
    x = Inches(0.6 + i * 4.2)
    # Title box
    tbox = slide.shapes.add_shape(1, x, Inches(1.3), Inches(3.8), Inches(1.0))
    tbox.fill.solid()
    tbox.fill.fore_color.rgb = PURPLE_MID
    tbox.line.fill.background()
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(16)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.color.rgb = WHITE
    tbox.text_frame.paragraphs[0].font.name = "Calibri"
    tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Description box
    dbox = slide.shapes.add_shape(1, x, Inches(2.4), Inches(3.8), Inches(2.5))
    dbox.fill.solid()
    dbox.fill.fore_color.rgb = RGBColor(0xF8, 0xF4, 0xFC)
    dbox.line.color.rgb = PURPLE_LIGHT
    dbox.text_frame.paragraphs[0].text = desc
    dbox.text_frame.paragraphs[0].font.size = Pt(13)
    dbox.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    dbox.text_frame.paragraphs[0].font.name = "Calibri"
    dbox.text_frame.word_wrap = True

# Three style rendering boxes at bottom
styles_info = [
    ("Flat Style", "Arial Bold\nDominant color fill\nWhite stroke + shadow", PURPLE_DARK),
    ("Cartoon Style", "Comic Sans Bold\nWhite fill\n4px colored outline", PURPLE_MID),
    ("Watercolor Style", "Segoe Script\nHandwritten font\nGaussian glow effect", PURPLE_LIGHT),
]

for i, (title, desc, color) in enumerate(styles_info):
    x = Inches(0.6 + i * 4.2)
    sbox = slide.shapes.add_shape(1, x, Inches(5.3), Inches(3.8), Inches(1.8))
    sbox.fill.solid()
    sbox.fill.fore_color.rgb = color
    sbox.line.fill.background()
    
    stf = sbox.text_frame
    stf.word_wrap = True
    stf.paragraphs[0].text = title
    stf.paragraphs[0].font.size = Pt(16)
    stf.paragraphs[0].font.bold = True
    stf.paragraphs[0].font.color.rgb = WHITE
    stf.paragraphs[0].font.name = "Calibri"
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p2 = stf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0xE8, 0xE0, 0xF0)
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

# ============================================================
# Slide 6: Three-Style Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_purple_header(slide, "Results: Three-Style Text Overlay Comparison")

img_path = REPORT_FIG_DIR / "text_overlay_3styles_animal_1.png"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(1.5), Inches(1.3), Inches(10.3), Inches(3.5))

desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.5), Inches(2))
dtf = desc_box.text_frame
dtf.word_wrap = True
add_first_bullet(dtf, "Same animal sticker with three different text rendering styles:", size=18, bold=True)
add_bullet_text(dtf, "Flat: Dominant color text with white stroke and drop shadow", size=16, level=1)
add_bullet_text(dtf, "Cartoon: White text with thick colored outline (comic book style)", size=16, level=1)
add_bullet_text(dtf, "Watercolor: Handwritten font with soft Gaussian glow", size=16, level=1)

# ============================================================
# Slide 7: Multi-Animal Grid
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_purple_header(slide, "Results: Batch Text Overlay (3 Animals × 3 Styles)")

img_path = REPORT_FIG_DIR / "text_overlay_grid.png"
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(2), Inches(1.2), Inches(9.3), Inches(5))

stat_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8))
stf = stat_box.text_frame
stf.word_wrap = True
p = stf.paragraphs[0]
p.text = "Total output: 37 stickers × 3 styles = 111 final sticker images with VLM-generated captions"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = PURPLE_DARK
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# ============================================================
# Slide 8: Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_purple_header(slide, "Summary: Textual Generate")

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True

add_first_bullet(tf, "Problem Solved", size=22, bold=True, color=PURPLE_DARK)
add_bullet_text(tf, "SD 1.5 cannot generate readable text → Post-processing with VLM + PIL", size=18, level=1)
add_bullet_text(tf, "", size=10)

add_bullet_text(tf, "VLM Caption Generation", size=22, bold=True, color=PURPLE_DARK)
add_bullet_text(tf, "GPT-4o-mini automatically generates adjective + noun captions", size=18, level=1)
add_bullet_text(tf, "Diversity constraints prevent repetitive outputs", size=18, level=1)
add_bullet_text(tf, "", size=10)

add_bullet_text(tf, "Style-Aware Rendering", size=22, bold=True, color=PURPLE_DARK)
add_bullet_text(tf, "Dominant color extraction + local brightness sampling + contrast adjustment", size=18, level=1)
add_bullet_text(tf, "Three distinct styles: Flat / Cartoon / Watercolor", size=18, level=1)
add_bullet_text(tf, "", size=10)

add_bullet_text(tf, "Output", size=22, bold=True, color=PURPLE_DARK)
add_bullet_text(tf, "37 stickers × 3 styles = 111 final images with consistent quality", size=18, level=1)

# ============================================================
# Save
# ============================================================
prs.save(str(OUTPUT_PPT))
print(f"PPT saved to: {OUTPUT_PPT}")
print("Done!")
